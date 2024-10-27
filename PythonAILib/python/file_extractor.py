from magika import Magika  # type: ignore
import chardet
import excel_util

def identify_type(filename):
    m = Magika()
    # ファイルのbyte列を取得
    with open(filename, "rb") as f:
        byte_data = f.read()

    # ファイルの種類を判定
    res = m.identify_bytes(byte_data)
    # エンコーディング判定
    encoding_dic = chardet.detect(byte_data)
    encoding = encoding_dic["encoding"]
    if encoding is None or encoding == "UTF-8":
        encoding = "utf-8"
    elif encoding == "SHIFT_JIS":
        encoding = "cp932"
           
    return res, encoding

def get_mime_type(filename):
    res, encoding = identify_type(filename)
    return res.output.mime_type

def extract_text_from_file(filename):
    res, encoding = identify_type(filename)
    print(res.output.mime_type)
    
    if res.output.mime_type.startswith("text/"):
        return process_text(filename, res, encoding)

    # application/pdf
    elif res.output.mime_type == "application/pdf":
        return process_pdf(filename)
        
    # application/vnd.openxmlformats-officedocument.spreadsheetml.sheet
    elif res.output.mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        return excel_util.extract_text_from_sheet(filename)
        
    # application/vnd.openxmlformats-officedocument.wordprocessingml.document
    elif res.output.mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return process_docx(filename)
        
    # application/vnd.openxmlformats-officedocument.presentationml.presentation
    elif res.output.mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return process_pptx(filename)
    else:
        print("Unsupported file type: " + res.output.mime_type)



# application/pdfのファイルを読み込んで文字列として返す関数
def process_pdf(filename):
    from pdfminer.high_level import extract_text
    text = extract_text(filename)
    return text


# application/vnd.openxmlformats-officedocument.wordprocessingml.documentのファイルを読み込んで文字列として返す関数
def process_docx(filename):
    import docx
    from io import StringIO
    # 出力用のストリームを作成
    output = StringIO()
    doc = docx.Document(filename)
    for para in doc.paragraphs:
        output.write(para.text)
        output.write("\n")
        
    return output.getvalue()

# application/vnd.openxmlformats-officedocument.presentationml.presentationのファイルを読み込んで文字列として返す関数
def process_pptx(filename):
    import pptx
    from io import StringIO
    # 出力用のストリームを作成
    output = StringIO()
    prs = pptx.Presentation(filename)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                output.write(shape.text)
                output.write("\n")
    
    return output.getvalue()

# text/*のファイルを読み込んで文字列として返す関数
def process_text(filename, res, encoding):
    result = ""
    if res.output.mime_type == "text/html":
        # text/htmlの場合
        from bs4 import BeautifulSoup
        # テキストを取得
        with open(filename, "rb") as f:
            text_data = f.read()
            soup = BeautifulSoup(text_data, "html.parser")
        result = soup.get_text()

    elif res.output.mime_type == "text/xml":
        # text/xmlの場合
        from bs4 import BeautifulSoup
        # テキストを取得
        with open(filename, "rb") as f:
            text_data = f.read()
            soup = BeautifulSoup(text_data, features="xml")
        result = soup.get_text()

    elif res.output.mime_type == "text/markdown":
        # markdownの場合
        from bs4 import BeautifulSoup
        from markdown import markdown
        # テキストを取得
        with open(filename, "r" ,encoding=encoding) as f:
            text_data = f.read()
            md = markdown(text_data)
            soup = BeautifulSoup(md, "html.parser")
        result = soup.get_text()
    else:
        # その他のtext/*の場合
        with open(filename, "r", encoding=encoding) as f:
            result = f.read()
        
    return result

